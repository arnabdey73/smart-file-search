"""
Content extraction module.
Handles extracting text content from various file formats.
"""

import os
import asyncio
import zipfile
from pathlib import Path
from typing import List, Optional
import logging

logger = logging.getLogger(__name__)

# Optional imports with graceful fallbacks
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from openpyxl import load_workbook
    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from bs4 import BeautifulSoup
    HAS_BS4 = True
except ImportError:
    HAS_BS4 = False


class ContentExtractor:
    """Extracts text content from various file formats."""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 100):
        """
        Initialize the content extractor.
        
        Args:
            chunk_size: Target size for text chunks
            chunk_overlap: Overlap between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    async def extract_content(self, file_path: str) -> List[str]:
        """
        Extract text content from a file and return as chunks.
        
        Args:
            file_path: Path to the file to extract
            
        Returns:
            List of text chunks
        """
        try:
            ext = Path(file_path).suffix.lower()
            
            # Route to appropriate extractor
            if ext == '.pdf':
                content = await self._extract_pdf(file_path)
            elif ext in ['.docx', '.doc']:
                content = await self._extract_docx(file_path)
            elif ext in ['.pptx', '.ppt']:
                content = await self._extract_pptx(file_path)
            elif ext in ['.xlsx', '.xls']:
                content = await self._extract_excel(file_path)
            elif ext == '.csv':
                content = await self._extract_csv(file_path)
            elif ext in ['.txt', '.md', '.py', '.js', '.ts', '.html', '.css', '.json', '.xml', '.yaml', '.yml', '.sql']:
                content = await self._extract_text(file_path)
            elif ext == '.zip':
                content = await self._extract_zip(file_path)
            else:
                logger.debug(f"Unsupported file type: {ext}")
                return []
            
            if not content:
                return []
            
            # Split into chunks
            return self._chunk_text(content)
            
        except Exception as e:
            logger.error(f"Error extracting content from {file_path}: {e}")
            return []
    
    async def _extract_pdf(self, file_path: str) -> Optional[str]:
        """Extract text from PDF file."""
        if not HAS_PDF:
            logger.warning("PyPDF2 not installed, skipping PDF extraction")
            return None
        
        try:
            text_parts = []
            
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                
                for page in reader.pages:
                    text = page.extract_text()
                    if text.strip():
                        text_parts.append(text.strip())
            
            return '\n\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting PDF {file_path}: {e}")
            return None
    
    async def _extract_docx(self, file_path: str) -> Optional[str]:
        """Extract text from DOCX file."""
        if not HAS_DOCX:
            logger.warning("python-docx not installed, skipping DOCX extraction")
            return None
        
        try:
            doc = Document(file_path)
            text_parts = []
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    text_parts.append(text)
            
            return '\n\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting DOCX {file_path}: {e}")
            return None
    
    async def _extract_pptx(self, file_path: str) -> Optional[str]:
        """Extract text from PPTX file."""
        if not HAS_PPTX:
            logger.warning("python-pptx not installed, skipping PPTX extraction")
            return None
        
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
            
            return '\n\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting PPTX {file_path}: {e}")
            return None
    
    async def _extract_excel(self, file_path: str) -> Optional[str]:
        """Extract text from Excel file."""
        if not HAS_EXCEL:
            logger.warning("openpyxl not installed, skipping Excel extraction")
            return None
        
        try:
            workbook = load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Add sheet name
                text_parts.append(f"Sheet: {sheet_name}")
                
                # Extract cell values
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell_value in row:
                        if cell_value is not None:
                            row_text.append(str(cell_value))
                    
                    if row_text:
                        text_parts.append('\t'.join(row_text))
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting Excel {file_path}: {e}")
            return None
    
    async def _extract_csv(self, file_path: str) -> Optional[str]:
        """Extract text from CSV file."""
        try:
            import csv
            
            text_parts = []
            
            with open(file_path, 'r', encoding='utf-8', newline='') as file:
                # Try to detect delimiter
                sample = file.read(1024)
                file.seek(0)
                
                sniffer = csv.Sniffer()
                delimiter = sniffer.sniff(sample).delimiter
                
                reader = csv.reader(file, delimiter=delimiter)
                
                for row in reader:
                    if row:
                        text_parts.append('\t'.join(str(cell) for cell in row))
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting CSV {file_path}: {e}")
            return None
    
    async def _extract_text(self, file_path: str) -> Optional[str]:
        """Extract text from plain text files."""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                    
                    # Basic cleanup for HTML/XML
                    if file_path.lower().endswith(('.html', '.htm', '.xml')):
                        content = self._clean_html(content)
                    
                    return content
                    
                except UnicodeDecodeError:
                    continue
            
            logger.warning(f"Could not decode text file {file_path}")
            return None
            
        except Exception as e:
            logger.error(f"Error extracting text {file_path}: {e}")
            return None
    
    def _clean_html(self, content: str) -> str:
        """Clean HTML content using BeautifulSoup if available."""
        if not HAS_BS4:
            # Basic HTML tag removal
            import re
            content = re.sub(r'<[^>]+>', '', content)
            return content
        
        try:
            soup = BeautifulSoup(content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text and clean up whitespace
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            return text
            
        except Exception:
            # Fallback to basic tag removal
            import re
            return re.sub(r'<[^>]+>', '', content)
    
    async def _extract_zip(self, file_path: str) -> Optional[str]:
        """Extract text from ZIP file (file names and small text files)."""
        try:
            text_parts = []
            
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Add file listing
                text_parts.append("ZIP Contents:")
                for info in zip_file.filelist:
                    text_parts.append(f"  {info.filename} ({info.file_size} bytes)")
                
                # Extract content from small text files
                for info in zip_file.filelist:
                    if info.file_size > 10000:  # Skip large files
                        continue
                    
                    if info.filename.lower().endswith(('.txt', '.md', '.json', '.xml', '.csv')):
                        try:
                            with zip_file.open(info) as inner_file:
                                content = inner_file.read().decode('utf-8')
                                text_parts.append(f"\n--- {info.filename} ---\n{content}")
                        except Exception:
                            continue
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error extracting ZIP {file_path}: {e}")
            return None
    
    def _chunk_text(self, text: str) -> List[str]:
        """Split text into overlapping chunks."""
        if len(text) <= self.chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + self.chunk_size
            
            # Try to break at word boundary
            if end < len(text):
                # Look for word boundary
                for i in range(end, max(start + self.chunk_size // 2, end - 100), -1):
                    if text[i] in ' \n\t.!?':
                        end = i + 1
                        break
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move start position with overlap
            start = end - self.chunk_overlap
            if start >= len(text):
                break
        
        return chunks
